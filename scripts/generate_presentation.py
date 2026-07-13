#!/usr/bin/env python3
import os
import sys

# Define absolute paths and insert virtual environment site-packages
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, os.path.join(ROOT, ".venv", "lib", "python3.12", "site-packages"))
sys.path.insert(0, os.path.join(ROOT, "src"))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Output Path
OUTPUT_DIR = os.path.join(ROOT, "reports", "ppt_presentation")
os.makedirs(OUTPUT_DIR, exist_ok=True)
OUTPUT_PATH = os.path.join(OUTPUT_DIR, "pgi_presentation.pptx")

# Design system constants (subtle Big 4 styling)
COLOR_BG_WHITE = RGBColor(255, 255, 255)
COLOR_BG_LIGHT = RGBColor(248, 250, 252)    # Slate 50 (light background panels)
COLOR_NAVY = RGBColor(15, 30, 54)          # Deep Corporate Navy (#0F1E36)
COLOR_CHARCOAL = RGBColor(30, 41, 59)      # Slate 800 (primary body text)
COLOR_SLATE = RGBColor(100, 116, 139)      # Slate 500 (secondary text)
COLOR_BORDER = RGBColor(226, 232, 240)     # Slate 200 (subtle separators)
COLOR_ACCENT_AMBER = RGBColor(180, 83, 9)  # Amber 700 (selective metrics)
COLOR_ACCENT_TEAL = RGBColor(15, 118, 110) # Teal 700 (selective metrics)
COLOR_TAKEAWAY_BG = RGBColor(254, 243, 199) # Light Amber 100 for takeaways
COLOR_TAKEAWAY_TXT = RGBColor(146, 64, 14) # Amber 800

FONT_NAME = "Calibri"

# Slide Dimensions (16:9 Widescreen)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Layout coordinates
LEFT_MARGIN = Inches(0.8)
RIGHT_MARGIN = Inches(12.533)
USABLE_WIDTH = RIGHT_MARGIN - LEFT_MARGIN
BODY_TOP = Inches(1.3)
BODY_BOTTOM = Inches(6.0)
BODY_HEIGHT = BODY_BOTTOM - BODY_TOP
TAKEAWAY_TOP = Inches(6.1)
TAKEAWAY_HEIGHT = Inches(0.7)

# Helper functions for slide building
def create_base_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs

def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_header(slide, title, category="RESEARCH METHODOLOGY"):
    # Kicker / Category Tracker
    kicker_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.35), USABLE_WIDTH, Inches(0.3))
    tf_k = kicker_box.text_frame
    tf_k.word_wrap = True
    tf_k.margin_left = tf_k.margin_top = tf_k.margin_right = tf_k.margin_bottom = 0
    p_k = tf_k.paragraphs[0]
    p_k.text = category.upper()
    p_k.font.name = FONT_NAME
    p_k.font.size = Pt(9)
    p_k.font.bold = True
    p_k.font.color.rgb = COLOR_SLATE
    
    # Title
    title_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(0.55), USABLE_WIDTH, Inches(0.5))
    tf_t = title_box.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
    p_t = tf_t.paragraphs[0]
    p_t.text = title
    p_t.font.name = FONT_NAME
    p_t.font.size = Pt(22)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_NAVY
    
    # Divider Line
    connector = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, LEFT_MARGIN, Inches(1.1), USABLE_WIDTH, Inches(0.02)
    )
    connector.fill.solid()
    connector.fill.fore_color.rgb = COLOR_BORDER
    connector.line.fill.background() # No border line

def add_slide_footer(slide, current_slide, total_slides):
    # Footer text
    footer_box = slide.shapes.add_textbox(LEFT_MARGIN, Inches(7.0), Inches(9.0), Inches(0.3))
    tf = footer_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = "Procurement Governance Intelligence (PGI)  ·  MSc Dissertation Defense"
    p.font.name = FONT_NAME
    p.font.size = Pt(9.5)
    p.font.color.rgb = COLOR_SLATE
    
    # Page number
    page_box = slide.shapes.add_textbox(Inches(11.533), Inches(7.0), Inches(1.0), Inches(0.3))
    tf_p = page_box.text_frame
    tf_p.word_wrap = True
    tf_p.margin_left = tf_p.margin_top = tf_p.margin_right = tf_p.margin_bottom = 0
    p_p = tf_p.paragraphs[0]
    p_p.alignment = PP_ALIGN.RIGHT
    p_p.text = f"{current_slide} / {total_slides}"
    p_p.font.name = FONT_NAME
    p_p.font.size = Pt(9.5)
    p_p.font.color.rgb = COLOR_SLATE

def add_bullet_points(slide, left, top, width, height, points, font_size=13):
    """
    points: list of tuples (bullet_type, text_content).
            bullet_type: '0' for primary bullet, '1' for secondary indented bullet.
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    for idx, (b_type, text) in enumerate(points):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        p.text = text
        p.font.name = FONT_NAME
        
        # Big 4 bullet indentation styling
        if b_type == 0:
            p.level = 0
            p.space_after = Pt(8)
            p.font.size = Pt(font_size)
            p.font.color.rgb = COLOR_CHARCOAL
            # Bold leading words if they contain a colon
            if ":" in text:
                parts = text.split(":", 1)
                p.text = ""
                run_bold = p.add_run()
                run_bold.text = parts[0] + ":"
                run_bold.font.bold = True
                run_bold.font.color.rgb = COLOR_NAVY
                
                run_normal = p.add_run()
                run_normal.text = parts[1]
                run_normal.font.bold = False
        else:
            p.level = 1
            p.space_after = Pt(4)
            p.font.size = Pt(font_size - 1)
            p.font.color.rgb = COLOR_SLATE
            if ":" in text:
                parts = text.split(":", 1)
                p.text = ""
                run_bold = p.add_run()
                run_bold.text = parts[0] + ":"
                run_bold.font.bold = True
                run_bold.font.color.rgb = COLOR_NAVY
                
                run_normal = p.add_run()
                run_normal.text = parts[1]
                run_normal.font.bold = False

def add_card_shape(slide, left, top, width, height, bg_color=COLOR_BG_LIGHT, border_color=COLOR_BORDER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
    else:
        card.line.fill.background()
    return card

def add_key_takeaway(slide, text):
    # Background Box
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, LEFT_MARGIN, TAKEAWAY_TOP, USABLE_WIDTH, TAKEAWAY_HEIGHT
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLOR_TAKEAWAY_BG
    rect.line.color.rgb = COLOR_BORDER
    rect.line.width = Pt(0.75)
    
    # Text Frame
    tb = slide.shapes.add_textbox(
        LEFT_MARGIN + Inches(0.15), TAKEAWAY_TOP + Inches(0.08), USABLE_WIDTH - Inches(0.3), TAKEAWAY_HEIGHT - Inches(0.16)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = ""
    
    run_prefix = p.add_run()
    run_prefix.text = "Key Takeaway: "
    run_prefix.font.name = FONT_NAME
    run_prefix.font.size = Pt(11)
    run_prefix.font.bold = True
    run_prefix.font.color.rgb = COLOR_TAKEAWAY_TXT
    
    run_text = p.add_run()
    run_text.text = text
    run_text.font.name = FONT_NAME
    run_text.font.size = Pt(11)
    run_text.font.bold = False
    run_text.font.color.rgb = COLOR_TAKEAWAY_TXT

def add_kpi_callout(slide, left, top, width, height, value, label, value_color=COLOR_NAVY):
    # Background card panel
    add_card_shape(slide, left, top, width, height, bg_color=COLOR_BG_LIGHT, border_color=COLOR_BORDER)
    
    # Text Container
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), height - Inches(0.3))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    # Large Number
    p_val = tf.paragraphs[0]
    p_val.text = value
    p_val.alignment = PP_ALIGN.CENTER
    p_val.font.name = FONT_NAME
    p_val.font.size = Pt(36)
    p_val.font.bold = True
    p_val.font.color.rgb = value_color
    p_val.space_after = Pt(2)
    
    # Small Label
    p_lbl = tf.add_paragraph()
    p_lbl.text = label
    p_lbl.alignment = PP_ALIGN.CENTER
    p_lbl.font.name = FONT_NAME
    p_lbl.font.size = Pt(10)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = COLOR_SLATE

def add_table(slide, left, top, width, height, headers, rows):
    rows_cnt = len(rows) + 1
    cols_cnt = len(headers)
    
    table_shape = slide.shapes.add_table(rows_cnt, cols_cnt, left, top, width, height)
    table = table_shape.table

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_NAVY
        # Alignment & font
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
        p.font.name = FONT_NAME
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = COLOR_BG_WHITE

    for row_idx, row_data in enumerate(rows):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(val)
            cell.fill.solid()
            # Zebra striping
            bg = COLOR_BG_LIGHT if row_idx % 2 == 0 else COLOR_BG_WHITE
            cell.fill.fore_color.rgb = bg
            
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
            p.font.name = FONT_NAME
            p.font.size = Pt(9.5)
            p.font.color.rgb = COLOR_CHARCOAL
            p.font.bold = False


# Presentation generation
def generate_deck():
    prs = create_base_presentation()
    total_slides = 13
    
    # Slide 1: title slide (cover slide)
    slide_layout = prs.slide_layouts[6] # Blank slide
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, COLOR_BG_WHITE)
    
    # Clean Big 4 cover styling: Vertical navy accent bar on left edge
    navy_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), SLIDE_HEIGHT)
    navy_bar.fill.solid()
    navy_bar.fill.fore_color.rgb = COLOR_NAVY
    navy_bar.line.fill.background()
    
    # Main Titles Box
    title_tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(3.0))
    tf = title_tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p_kicker = tf.paragraphs[0]
    p_kicker.text = "MSC IN ARTIFICIAL INTELLIGENCE  ·  DISSERTATION DEFENSE"
    p_kicker.font.name = FONT_NAME
    p_kicker.font.size = Pt(11)
    p_kicker.font.bold = True
    p_kicker.font.color.rgb = COLOR_SLATE
    p_kicker.space_after = Pt(12)
    
    p_main = tf.add_paragraph()
    p_main.text = "Procurement Governance Intelligence (PGI)"
    p_main.font.name = FONT_NAME
    p_main.font.size = Pt(40)
    p_main.font.bold = True
    p_main.font.color.rgb = COLOR_NAVY
    p_main.space_after = Pt(6)
    
    p_sub = tf.add_paragraph()
    p_sub.text = "A Network, Risk, and Price Analytics Framework for Italian Public Procurement"
    p_sub.font.name = FONT_NAME
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = COLOR_SLATE
    p_sub.space_after = Pt(48)
    
    p_details = tf.add_paragraph()
    p_details.text = "Submitted by: Shreyas Sudarshanam (MSc Candidate)  |  National College of Ireland\nSupervisor: Dissertation Committee  |  Date: June 2026"
    p_details.font.name = FONT_NAME
    p_details.font.size = Pt(13)
    p_details.font.bold = True
    p_details.font.color.rgb = COLOR_CHARCOAL
    
    # Add footer-like marker for cover page
    add_slide_footer(slide, 1, total_slides)
    
    # Slide 2: background & context
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, COLOR_BG_WHITE)
    add_slide_header(slide, "Background & Regulatory Governance Context", "Dissertation Context")
    add_slide_footer(slide, 2, total_slides)
    
    points_left = [
        (0, "Scale of Procurement: Public purchasing represents a substantial category of state activity across the EU, making it a primary locus for governance and compliance risks."),
        (0, "Worrying Competition Decline: The European Court of Auditors (ECA) has documented a long-term decline in competitive bids, characterized by rising single-bidder rates."),
        (0, "Relational Structures: Standard red-flag models assess contracts individually, ignoring the complex buyer-supplier transaction network.")
    ]
    points_right = [
        (0, "The Auditability Mandate: Decisions affecting public funds and private suppliers carry legal accountability requirements (EU Directive 2014/24/EU, GDPR Article 22)."),
        (0, "The XAI Constraint: Unexplained, opaque machine-learning risk scores are legally problematic and practically useless for national auditing bodies (like Italy's ANAC)."),
        (0, "Framework Philosophy: Prioritize model interpretability, deterministic data lineage, and local-first compute over accuracy-only black-box architectures.")
    ]
    
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.3), points_left)
    add_bullet_points(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.3), points_right)
    add_key_takeaway(slide, "Governance requirements make transparent decision logic a binding constraint, shaping every algorithm choice in the PGI framework.")

    # Slide 3: research pivots & data reality checks
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, COLOR_BG_WHITE)
    add_slide_header(slide, "Honest Reporting: Data Sparsity & Research Pivots", "Data Audit & Scope")
    add_slide_footer(slide, 3, total_slides)
    
    # Left text column
    points_left = [
        (0, "The Original Scope: The research was initially framed to model contract-delay durations and predict administrative bottlenecks."),
        (0, "The Data Sparsity Wall: A comprehensive data-quality audit of the Italian subset of the Global Public Procurement Dataset (GPPD) (~12.1M records, 2006-2021) revealed severe gaps:"),
        (1, "Contract Signature Date: 99.997% missing (only 346 records contain signature dates in the entire database)."),
        (1, "Award Decision Date: 98.30% missing."),
        (0, "The Research Pivot: Delay modeling on observed timelines was fundamentally blocked. redrafted scope toward Network, Risk, and Price Intelligence.")
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(6.8), Inches(4.3), points_left)
    
    # Right KPIs panel
    add_kpi_callout(slide, Inches(8.0), Inches(1.5), Inches(4.3), Inches(1.3), "99.997%", "SIGNATURE DATE NULL RATE", COLOR_ACCENT_AMBER)
    add_kpi_callout(slide, Inches(8.0), Inches(3.0), Inches(4.3), Inches(1.3), "98.30%", "AWARD DECISION DATE NULL RATE", COLOR_ACCENT_AMBER)
    add_kpi_callout(slide, Inches(8.0), Inches(4.5), Inches(4.3), Inches(1.3), "1.6%", "LOT-LEVEL BID COUNT COVERAGE", COLOR_CHARCOAL)
    
    add_key_takeaway(slide, "Data-driven honest reporting required pivoting away from delay prediction and highlights the coverage limits of standard procurement records.")

    # Slide 4: system architecture & framework pillars
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, COLOR_BG_WHITE)
    add_slide_header(slide, "SQL-First, Local-First System Architecture", "System Design")
    add_slide_footer(slide, 4, total_slides)
    
    points_left = [
        (0, "SQL-First Feature Store: Transformations are built declaratively using DuckDB and Polars. This guarantees that all model features have an inspectable SQL lineage."),
        (0, "Local-First Compute: Designed to run locally on resource-constrained hardware to keep data secure. A lightweight AWS migration path (S3 + Athena) is mapped but not mandated."),
        (0, "Governance Layer: Enforces strict data contracts, temporal split boundaries, and the automated leakage-audit suite before allowing any downstream model scoring.")
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.3), points_left)
    
    # Right Side: Framework Pillars
    add_card_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.3), COLOR_BG_LIGHT, COLOR_BORDER)
    
    tbl_tb = slide.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.1), Inches(3.9))
    tf = tbl_tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    
    p_ph = tf.paragraphs[0]
    p_ph.text = "THE THREE ANALYTIC PILLARS OF PGI"
    p_ph.font.name = FONT_NAME
    p_ph.font.size = Pt(11)
    p_ph.font.bold = True
    p_ph.font.color.rgb = COLOR_NAVY
    p_ph.space_after = Pt(12)
    
    p_p1 = tf.add_paragraph()
    p_p1.text = "1. Supplier Network Intelligence (RQ1)\nModels relational structures as a bipartite graph, applying Louvain modularity, Brandes centrality, and network resilience tests."
    p_p1.font.name = FONT_NAME
    p_p1.font.size = Pt(12)
    p_p1.font.color.rgb = COLOR_CHARCOAL
    p_p1.space_after = Pt(10)
    
    p_p2 = tf.add_paragraph()
    p_p2.text = "2. Governance Risk Intelligence (RQ2)\nFrames risk classification under a strict temporal split, leveraging white-box tree ensembles (XGBoost, EBM) and a 20+ check leakage audit."
    p_p2.font.name = FONT_NAME
    p_p2.font.size = Pt(12)
    p_p2.font.color.rgb = COLOR_CHARCOAL
    p_p2.space_after = Pt(10)
    
    p_p3 = tf.add_paragraph()
    p_p3.text = "3. Category-Restricted Price Intelligence (RQ3)\nEstablishes price baselines within CPV families using robust regressors and identifies anomalous outliers via an unsupervised consensus ensemble."
    p_p3.font.name = FONT_NAME
    p_p3.font.size = Pt(12)
    p_p3.font.color.rgb = COLOR_CHARCOAL
    
    add_key_takeaway(slide, "The modular architecture integrates Bipartite Networks (RQ1), Tree Ensembles (RQ2), and Outlier Ensembles (RQ3) into a single unified dashboard.")

    # Slide 5: RQ1, supplier network intel
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, COLOR_BG_WHITE)
    add_slide_header(slide, "RQ1: Bipartite Network Construction & Centrality", "Research Question 1")
    add_slide_footer(slide, 5, total_slides)
    
    points_left = [
        (0, "Graph Formulation: Constructed a bipartite buyer-supplier network using the training window (2014-2018). Nodes are buyers and suppliers; edges represent contracts weighted by award value."),
        (0, "Scale & Sizing: Contains 161,695 nodes and 296,306 edges. Features 18,357 unique buyers and 147,609 unique suppliers. Earlier data caps were discarded to avoid inflating metrics."),
        (0, "Centrality Concentration: Combined normalized degree and Brandes betweenness ($k=500$, seed $42$) to construct a composite importance score. Centrality is heavy-tailed: the top decile concentration is ~9x the mean.")
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.3), points_left)
    
    # Right Column: Network Structure Table
    add_kpi_callout(slide, Inches(6.8), Inches(1.5), Inches(2.7), Inches(1.3), "161.7k", "UNIQUE NODES IN GRAPH", COLOR_NAVY)
    add_kpi_callout(slide, Inches(9.8), Inches(1.5), Inches(2.7), Inches(1.3), "296.3k", "TOTAL NETWORK EDGES", COLOR_NAVY)
    
    table_headers = ["Network Segment", "Nodes Count", "Modularity (Q)", "ARI Stability"]
    table_rows = [
        ["Whole Graph", "161,695", "0.7258", "0.5573"],
        ["Giant Component", "159,893", "0.7137", "0.5462"],
        ["Isolated Hubs", "1,802", "0.0121", "N/A"]
    ]
    add_table(slide, Inches(6.8), Inches(3.0), Inches(5.7), Inches(1.4), table_headers, table_rows)
    
    add_key_takeaway(slide, "The Italian procurement network is dominated by a giant component (98.89% of nodes) with non-trivial community structure (Modularity Q = 0.726).")

    # Slide 6: RQ1, modularity null & topological validation
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, COLOR_BG_WHITE)
    add_slide_header(slide, "RQ1: Confronting the Modularity Null & Validation", "Research Question 1")
    add_slide_footer(slide, 6, total_slides)
    
    points_left = [
        (0, "Negative Modularity Null: Observed modularity $Q=0.726$ is *lower* than the degree-preserving null model mean $Q_{\\text{null}}=0.774$ ($z = -28.5$, $p=1.0$)."),
        (1, "Modularity magnitude alone does *not* prove significant structure. Bipartite networks naturally yield high modularity due to random fluctuation (Guimerà 2004)."),
        (0, "Topological Validation: Validity is instead established via conductance and cluster stability rather than a single modularity score."),
        (1, "Conductance: Giant component mean conductance is 0.281, with 64.2% of communities below the 0.30 cut threshold (genuine sparse cuts).")
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.3), points_left)
    
    points_right = [
        (0, "Louvain Stability: Adjusted Rand Index (ARI) over 5 random seeds yields a mean of 0.557, representing moderate stability."),
        (1, "Stability limitation is reported transparently (under the 0.70 target) since instability concentrates in small, peripheral communities."),
        (0, "Buyer Resilience: Dispersion of buyer-level resilience has standard deviation $\\sigma = 0.270 > 0.15$ (gate met)."),
        (1, "Indicates public buyers are *not* uniformly resilient; some are highly vulnerable to single-supplier dependencies.")
    ]
    add_bullet_points(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.3), points_right)
    
    add_key_takeaway(slide, "Modularity significance was negative, but giant component conductance (mean 0.281) and resilience variance (0.270) establish topological validity.")

    # Slide 7: RQ2, governance risk & leakage audit
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, COLOR_BG_WHITE)
    add_slide_header(slide, "RQ2: Supervised Risk Modeling & Debiasing", "Research Question 2")
    add_slide_footer(slide, 7, total_slides)
    
    points_left = [
        (0, "Strict Temporal Split: Evaluation is structured to avoid future leakage: Train 2014-2018 ($n=471,136$), Validation 2019 ($n=120,205$), Test 2020 ($n=123,340$)."),
        (0, "Target Label Formulation: The governance-risk label is a binary flag aggregating sole-sourcing procedure type, low bid counts, and extreme buyer dependency ratio."),
        (0, "Algorithmic Debiasing: The label rule was modified to **exclude contract amendments** because a data quality audit revealed amendments were recording artifacts restricted to region ITH5."),
        (1, "Removing amendments prevents geographic bias under GDPR and EU AI Act principles.")
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.3), points_left)
    
    points_right = [
        (0, "Leakage-Audit Gate: Implemented an automated 20+ check suite evaluating target leakage, data leakage, and feature redundancy."),
        (0, "Synthetic-Fixture Control: synthetic check fixtures verify the audit script's sensitivity by embedding deliberate leaks and asserting failures."),
        (0, "No Leakage, No Model: Upstream model parameters are rejected unless the audit passes with zero warnings, satisfying trustworthy AI criteria.")
    ]
    add_bullet_points(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.3), points_right)
    
    add_key_takeaway(slide, "The target was debiased by removing regional amendment artifacts, and models are temporal-holdout gated by a strict 20+ leakage audit.")

    # Slide 8: RQ2, value-only feature & documented nulls
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, COLOR_BG_WHITE)
    add_slide_header(slide, "RQ2: Single Feature Architecture & Documented Nulls", "Research Question 2")
    add_slide_footer(slide, 8, total_slides)
    
    points_left = [
        (0, "The Value-Only Model: The leakage audit rejected seven candidate features. The only non-leaky feature retained was `contract_value_log`."),
        (0, "White-Box Classifiers: Trained three interpretable models on the single feature: Explainable Boosting Machine (EBM baseline), Logistic Regression, and constrained XGBoost (depth 3, headline model)."),
        (0, "SHAP Skip: Since the feature space is 1D, SHAP post-hoc attribution is degenerate and skipped, relying on Partial Dependence instead.")
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.3), points_left)
    
    # Right Side: Documented Nulls Table
    add_card_shape(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.3), COLOR_BG_LIGHT, COLOR_BORDER)
    
    null_tb = slide.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.1), Inches(3.9))
    tf_n = null_tb.text_frame
    tf_n.word_wrap = True
    tf_n.margin_left = tf_n.margin_top = tf_n.margin_right = tf_n.margin_bottom = 0
    
    p_nh = tf_n.paragraphs[0]
    p_nh.text = "REJECTED FEATURES (DOCUMENTED NULL RESULTS)"
    p_nh.font.name = FONT_NAME
    p_nh.font.size = Pt(11)
    p_nh.font.bold = True
    p_nh.font.color.rgb = COLOR_NAVY
    p_nh.space_after = Pt(12)
    
    p_n1 = tf_n.add_paragraph()
    p_n1.text = "• Supplier Centrality: Structural position adds no predictive signal; centrality alone yields test AUC of 0.490."
    p_n1.font.name = FONT_NAME
    p_n1.font.size = Pt(11.5)
    p_n1.font.color.rgb = COLOR_CHARCOAL
    p_n1.space_after = Pt(8)
    
    p_n2 = tf_n.add_paragraph()
    p_n2.text = "• Buyer Concentration (HHI): Bounded HHI exhibits target leakage due to high correlation (|r| = 0.77) with dependency ratio (a label input)."
    p_n2.font.name = FONT_NAME
    p_n2.font.size = Pt(11.5)
    p_n2.font.color.rgb = COLOR_CHARCOAL
    p_n2.space_after = Pt(8)
    
    p_n3 = tf_n.add_paragraph()
    p_n3.text = "• Buyer Region: Strongly correlated (r = 0.966) with regional amendments recording artifact; marginal lift of +0.001 was rejected."
    p_n3.font.name = FONT_NAME
    p_n3.font.size = Pt(11.5)
    p_n3.font.color.rgb = COLOR_CHARCOAL
    p_n3.space_after = Pt(8)
    
    p_n4 = tf_n.add_paragraph()
    p_n4.text = "• Tender Lot Count: Size proxy redundant with contract value (r = 0.53); rejected to keep model strictly value-only."
    p_n4.font.name = FONT_NAME
    p_n4.font.size = Pt(11.5)
    p_n4.font.color.rgb = COLOR_CHARCOAL
    
    add_key_takeaway(slide, "The leakage audit reduced the feature space to a single non-leaky predictor, producing 4 documented null results.")

    # Slide 9: RQ2, classification performance & base-rate drift
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, COLOR_BG_WHITE)
    add_slide_header(slide, "RQ2: Performance & Real-World Non-Stationarity", "Research Question 2")
    add_slide_footer(slide, 9, total_slides)
    
    points_left = [
        (0, "Headline Performance: Constrained XGBoost achieves a threshold-free ROC-AUC of **0.826** and PR-AUC of **0.834** on the 2020 temporal test set."),
        (0, "Baseline Comparison: XGBoost edges out Logistic Regression (AUC 0.808) and Random Forest (EBM) (AUC 0.704). DeLong paired tests confirm significance (p < 0.001)."),
        (0, "Operating Point Drift: A threshold chosen on validation (0.356) yields a test F1 of 0.572 (Precision 0.412 / Recall 0.938). The default 0.5 threshold yields F1 of 0.766."),
        (0, "The COVID Base-Rate Shift: Shuffled 5-fold CV yields 0.710. Temporal forward-chaining CV falls to $0.612 \\pm 0.138$, exposing severe market shifts between folds (2019 validation rate ~10% vs. 2020 test rate ~37%).")
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(6.8), Inches(4.3), points_left)
    
    # Right KPIs panel
    add_kpi_callout(slide, Inches(8.0), Inches(1.5), Inches(4.3), Inches(1.3), "0.826", "2020 TEST SET HEADLINE ROC-AUC", COLOR_NAVY)
    add_kpi_callout(slide, Inches(8.0), Inches(3.0), Inches(4.3), Inches(1.3), "0.834", "2020 TEST SET PR-AUC", COLOR_NAVY)
    add_kpi_callout(slide, Inches(8.0), Inches(4.5), Inches(4.3), Inches(1.3), "0.612 +/- 0.138", "TEMPORAL CV SCORE (MARKET DRIFT)", COLOR_ACCENT_AMBER)
    
    add_key_takeaway(slide, "Partial dependence shows a non-monotonic risk curve (elevated at low and high contract values), explaining why tree models out-perform linear baselines.")

    # Slide 10: RQ3, category-restricted price intel
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, COLOR_BG_WHITE)
    add_slide_header(slide, "RQ3: Category-Restricted Price Intelligence", "Research Question 3")
    add_slide_footer(slide, 10, total_slides)
    
    points_left = [
        (0, "Category Restriction: Restricting analysis to three well-populated CPV families (Construction, Medical, and IT/Telecom) to ensure stable within-category pricing metrics (20,598 contracts)."),
        (0, "Expected Price Benchmark: Trained a Gradient Boosted Regressor to estimate expected contract prices."),
        (1, "Leakage Resolution: Excluding target leak `log_estimated_price` (r=0.63 with target) reduces $R^2$ from 0.91 to an honest **0.283** (RMSE 1.811)."),
        (0, "Consensus Ensemble: Integrated three unsupervised density/distance outliers: Isolation Forest, Local Outlier Factor (LOF), and Robust Covariance (Elliptic Envelope)."),
        (0, "Contamination Calibration: Calibrated anomaly threshold via contamination sweep to achieve a 5.08% consensus flag rate, identifying 1,046 outliers.")
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.3), points_left)
    
    # Right Column: KPI & Anomaly Rates Table
    add_kpi_callout(slide, Inches(6.8), Inches(1.5), Inches(2.7), Inches(1.3), "0.283", "HONEST BENCHMARK R² VALUE", COLOR_CHARCOAL)
    add_kpi_callout(slide, Inches(9.8), Inches(1.5), Inches(2.7), Inches(1.3), "5.08%", "CONSTITUENT CONSENSUS RATE", COLOR_ACCENT_TEAL)
    
    table_headers = ["Sector/Category", "Total Lots", "Consensus Outliers", "Outlier Rate"]
    table_rows = [
        ["Construction / Works", "9,850", "498", "5.06%"],
        ["Medical Supplies", "6,120", "312", "5.10%"],
        ["IT & Telecom", "4,628", "236", "5.10%"],
        ["Total Portfolio", "20,598", "1,046", "5.08%"]
    ]
    add_table(slide, Inches(6.8), Inches(3.0), Inches(5.7), Inches(1.4), table_headers, table_rows)
    
    add_key_takeaway(slide, "Target leakage was resolved to yield an honest price benchmark ($R^2 = 0.283$), with outliers flagged by a 5.08% calibrated consensus rate.")

    # Slide 11: cross-RQ integration layer
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, COLOR_BG_WHITE)
    add_slide_header(slide, "Cross-RQ Intelligence: Joining the Views", "Cross-RQ Layer")
    add_slide_footer(slide, 11, total_slides)
    
    points_left = [
        (0, "Standardized Joins: Standardized outputs from RQ1 (communities, centrality), RQ2 (governance risk probabilities), and RQ3 (price anomaly flags) were joined on shared keys ($n = 594,476$ contracts)."),
        (0, "Network Position vs. Risk: Pearson correlation is **r = -0.011** (Honest Null). Supplier centrality in the market is independent of administrative compliance risk."),
        (0, "Network Position vs. Price Premium: Uncomputed (reported as a pipeline gap)."),
        (0, "Governance Risk vs. Price Anomaly: Degenerate low coverage (only 8.3% paired contracts, observed $r=0.739$ is non-representative, headline reported as null).")
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.3), points_left)
    
    points_right = [
        (0, "Community Risk Analysis: Aggregated RQ2 risk and RQ3 anomalies to the network community level to locate compliance hotspots."),
        (0, "Statistical Verification: Conducted ANOVA and Chi-Square tests over 79 communities with at least 10 contracts."),
        (1, "ANOVA Result: High risk variation across communities ($F \\approx 20,961.59$, $p < 0.001$)."),
        (1, "Chi-Square Result: Outlier rates are significantly dependent on network communities ($p < 0.001$)."),
        (0, "Outcome: Proves governance risk and pricing anomalies cluster strongly within specific network groups.")
    ]
    add_bullet_points(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.3), points_right)
    
    add_key_takeaway(slide, "Network centrality and risk are independent, but ANOVA proves risk and pricing anomalies cluster within specific network communities ($p < 0.001$).")

    # Slide 12: business outcomes & policy value
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, COLOR_BG_WHITE)
    add_slide_header(slide, "PGI Business Outcomes & Policy Value", "Business Evaluation")
    add_slide_footer(slide, 12, total_slides)
    
    points_left = [
        (0, "Oversight Prioritization (ANAC / ECA): Shifts auditing from exhaustive random screening to a prioritized, explainable shortlist of high-risk outliers."),
        (1, "Combines community risk profiles (RQ1), value-driven compliance indicators (RQ2), and price variance screening (RQ3)."),
        (0, "Public Procurer Self-Assessment: Enables procurement agencies to run local self-assessments to detect supplier concentration and pricing deviation."),
        (0, "GDPR & EU AI Act Compliance: intrinsically transparent white-box models and explicit data lineage guarantee auditability.")
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(6.8), Inches(4.3), points_left)
    
    # Right Side: Actionable Deliverables Card
    add_card_shape(slide, Inches(8.0), Inches(1.5), Inches(4.3), Inches(4.3), COLOR_BG_LIGHT, COLOR_BORDER)
    
    deliv_tb = slide.shapes.add_textbox(Inches(8.2), Inches(1.7), Inches(3.9), Inches(3.9))
    tf_d = deliv_tb.text_frame
    tf_d.word_wrap = True
    tf_d.margin_left = tf_d.margin_top = tf_d.margin_right = tf_d.margin_bottom = 0
    
    p_dh = tf_d.paragraphs[0]
    p_dh.text = "ACTIONABLE POLICY OUTPUTS"
    p_dh.font.name = FONT_NAME
    p_dh.font.size = Pt(11)
    p_dh.font.bold = True
    p_dh.font.color.rgb = COLOR_NAVY
    p_dh.space_after = Pt(12)
    
    p_d1 = tf_d.add_paragraph()
    p_d1.text = "• Auditable Shortlists: White-box EBM & XGBoost outputs supply clear reasoning (Partial Dependence curves) for every risk ranking."
    p_d1.font.name = FONT_NAME
    p_d1.font.size = Pt(11)
    p_d1.font.color.rgb = COLOR_CHARCOAL
    p_d1.space_after = Pt(10)
    
    p_d2 = tf_d.add_paragraph()
    p_d2.text = "• Price Outlier Alerts: Unsupervised consensus anomalies highlight specific purchase categories experiencing abnormal price inflation."
    p_d2.font.name = FONT_NAME
    p_d2.font.size = Pt(11)
    p_d2.font.color.rgb = COLOR_CHARCOAL
    p_d2.space_after = Pt(10)
    
    p_d3 = tf_d.add_paragraph()
    p_d3.text = "• Community Risk Flags: ANOVA maps structural supplier groups that concentrate elevated risk, shielding buyers from localized monopolies."
    p_d3.font.name = FONT_NAME
    p_d3.font.size = Pt(11)
    p_d3.font.color.rgb = COLOR_CHARCOAL
    
    add_key_takeaway(slide, "PGI provides regulatory-compliant, explainable shortlists that help audit bodies prioritize and target oversight resources.")

    # Slide 13: limitations & future work
    slide = prs.slides.add_slide(slide_layout)
    apply_background(slide, COLOR_BG_WHITE)
    add_slide_header(slide, "Technical Limitations & Future Directions", "Limitations & Future Work")
    add_slide_footer(slide, 13, total_slides)
    
    points_left = [
        (0, "Database Sparsity Gaps: The raw dataset contains severe coverage limits. Bid counts are populated for only 1.6% of contracts, and CPV code categorization is limited (97.95% missing globally)."),
        (0, "Network Stability: Bipartite network Louvain community stability yields a moderate ARI of 0.557, with instabilities concentrated in small peripheral clusters."),
        (0, "Label Proxy Limits: The governance-risk label rule is a proxy for procedural deviation, not a direct measure of active corruption.")
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(5.6), Inches(4.3), points_left)
    
    points_right = [
        (0, "Future Extension - Scale: Expand pricing anomaly screening to additional CPV categories as data coverage improves."),
        (0, "Future Extension - Cloud: Test the scalability of the local-first design by migrating storage to AWS S3 and query compute to Athena."),
        (0, "Future Extension - Dynamic Networks: Model supplier networks temporally to trace how buyer dependencies shift over years.")
    ]
    add_bullet_points(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.3), points_right)
    
    add_key_takeaway(slide, "Future efforts focus on scaling CPV coverage and expanding local-first compute to cloud storage architectures.")
    
    prs.save(OUTPUT_PATH)
    print(f"Presentation saved to: {OUTPUT_PATH}")

if __name__ == "__main__":
    generate_deck()
